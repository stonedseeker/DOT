import os
import pandas as pd
import PyPDF2
from docx import Document
from pptx import Presentation
from typing import Dict, Any
import logging

class DocumentParser:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def parse_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Parse document based on file type"""
        try:
            if file_type.lower() == 'pdf':
                return self._parse_pdf(file_path)
            elif file_type.lower() == 'docx':
                return self._parse_docx(file_path)
            elif file_type.lower() == 'pptx':
                return self._parse_pptx(file_path)
            elif file_type.lower() == 'csv':
                return self._parse_csv(file_path)
            elif file_type.lower() in ['txt', 'md']:
                return self._parse_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            self.logger.error(f"Error parsing {file_type} file: {e}")
            raise
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF file"""
        text_content = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    text_content.append({
                        'page': page_num + 1,
                        'content': text.strip()
                    })
        
        return {
            'type': 'pdf',
            'total_pages': len(text_content),
            'content': text_content,
            'metadata': {'file_path': file_path}
        }
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX file"""
        doc = Document(file_path)
        paragraphs = []
        
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                paragraphs.append({
                    'paragraph': i + 1,
                    'content': paragraph.text.strip()
                })
        
        return {
            'type': 'docx',
            'total_paragraphs': len(paragraphs),
            'content': paragraphs,
            'metadata': {'file_path': file_path}
        }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX file"""
        presentation = Presentation(file_path)
        slides_content = []
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text.strip():
                        slide_text.append(shape.text.strip())
            
            if slide_text:
                slides_content.append({
                    'slide': slide_num + 1,
                    'content': ' '.join(slide_text)
                })
        
        return {
            'type': 'pptx',
            'total_slides': len(slides_content),
            'content': slides_content,
            'metadata': {'file_path': file_path}
        }
    
    def _parse_csv(self, file_path: str) -> Dict[str, Any]:
        """Parse CSV file"""
        df = pd.read_csv(file_path)
        
        # Convert to text representation
        csv_content = []
        
        # Add header information
        headers = list(df.columns)
        csv_content.append({
            'section': 'headers',
            'content': f"CSV Headers: {', '.join(headers)}"
        })
        
        # Add summary statistics
        csv_content.append({
            'section': 'summary',
            'content': f"Total rows: {len(df)}, Total columns: {len(df.columns)}"
        })
        
        # Add sample data (first 5 rows)
        sample_data = df.head().to_string()
        csv_content.append({
            'section': 'sample_data',
            'content': f"Sample data:\n{sample_data}"
        })
        
        return {
            'type': 'csv',
            'total_rows': len(df),
            'total_columns': len(df.columns),
            'content': csv_content,
            'metadata': {'file_path': file_path, 'columns': headers}
        }
    
    def _parse_text(self, file_path: str) -> Dict[str, Any]:
        """Parse text/markdown file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Split into chunks by paragraphs
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        
        text_content = []
        for i, paragraph in enumerate(paragraphs):
            text_content.append({
                'paragraph': i + 1,
                'content': paragraph
            })
        
        return {
            'type': 'text',
            'total_paragraphs': len(text_content),
            'content': text_content,
            'metadata': {'file_path': file_path}
        }